from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any
import csv
import html
import re


class ArtifactError(ValueError):
    pass


@dataclass
class GeneratedArtifact:
    content: bytes
    mime_type: str
    filename: str
    suffix: str


MIME_BY_FORMAT = {
    "txt": "text/plain; charset=utf-8",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf": "application/pdf",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "csv": "text/csv; charset=utf-8",
    "md": "text/markdown; charset=utf-8",
    "json": "application/json; charset=utf-8",
    "py": "text/x-python; charset=utf-8",
    "js": "text/javascript; charset=utf-8",
    "ts": "text/typescript; charset=utf-8",
    "html": "text/html; charset=utf-8",
    "css": "text/css; charset=utf-8",
    "xml": "application/xml; charset=utf-8",
    "yaml": "application/yaml; charset=utf-8",
    "yml": "application/yaml; charset=utf-8",
}


def _safe_stem(title: str) -> str:
    text = re.sub(r"[^\w\- ]+", "", str(title or "export"), flags=re.UNICODE).strip()
    text = re.sub(r"\s+", "-", text)[:48].strip("-")
    return text or "export"


def _attachment_names(message: dict[str, Any]) -> str:
    return ", ".join(str(item.get("original_name") or "") for item in message.get("attachments", []) if item.get("original_name"))


def _first_markdown_table(messages: list[dict[str, Any]]) -> list[list[str]] | None:
    separator = re.compile(r"^:?-{3,}:?$")
    for message in messages:
        lines = [line.strip() for line in str(message.get("content") or "").splitlines()]
        for index in range(len(lines) - 1):
            header_line = lines[index]
            separator_line = lines[index + 1]
            if "|" not in header_line or "|" not in separator_line:
                continue
            headers = [cell.strip() for cell in header_line.strip("|").split("|")]
            separators = [cell.strip() for cell in separator_line.strip("|").split("|")]
            if not headers or len(headers) != len(separators):
                continue
            if not all(separator.fullmatch(cell) for cell in separators):
                continue
            rows = [headers]
            for row_line in lines[index + 2 :]:
                if "|" not in row_line:
                    break
                cells = [cell.strip() for cell in row_line.strip("|").split("|")]
                if len(cells) != len(headers):
                    break
                rows.append(cells)
            return rows if len(rows) >= 2 else None
    return None


class ArtifactService:
    def generate(self, fmt: str, title: str, messages: list[dict[str, Any]]) -> GeneratedArtifact:
        fmt = str(fmt or "").lower().strip()
        if fmt not in MIME_BY_FORMAT:
            raise ArtifactError("Định dạng xuất không được hỗ trợ.")
        if not messages:
            raise ArtifactError("Không có nội dung để xuất.")
        stem = _safe_stem(title)
        if fmt in {"md", "json", "py", "js", "ts", "html", "css", "xml", "yaml", "yml"}:
            content = self._raw_text(messages, strip_fence=fmt in {"py", "js", "ts", "html", "css", "json", "xml", "yaml", "yml"})
        elif fmt == "csv":
            content = self._csv(messages)
        elif fmt == "txt":
            content = self._txt(title, messages)
        elif fmt == "docx":
            content = self._docx(title, messages)
        elif fmt == "pdf":
            content = self._pdf(title, messages)
        elif fmt == "pptx":
            content = self._pptx(title, messages)
        else:
            content = self._xlsx(title, messages)
        return GeneratedArtifact(content=content, mime_type=MIME_BY_FORMAT[fmt], filename=f"{stem}.{fmt}", suffix=f".{fmt}")


    def _raw_text(self, messages: list[dict[str, Any]], *, strip_fence: bool = False) -> bytes:
        content = "\n\n".join(
            str(item.get("content") or "").strip()
            for item in messages
            if str(item.get("content") or "").strip()
        ).strip()
        if strip_fence:
            match = re.fullmatch(r"```[^\n]*\n(.*?)\n```", content, flags=re.S)
            if match:
                content = match.group(1).strip()
        return content.encode("utf-8")

    def _csv(self, messages: list[dict[str, Any]]) -> bytes:
        table = _first_markdown_table(messages)
        import io
        text_buffer = io.StringIO()
        writer = csv.writer(text_buffer, lineterminator="\n")
        if table:
            writer.writerows(table)
        else:
            content = "\n\n".join(str(item.get("content") or "") for item in messages)
            for line in content.splitlines():
                writer.writerow([line])
        return text_buffer.getvalue().encode("utf-8-sig")

    def _txt(self, title: str, messages: list[dict[str, Any]]) -> bytes:
        lines = [str(title or "Cuộc trò chuyện"), "=" * max(12, len(str(title or ""))), ""]
        for item in messages:
            role = str(item.get("role") or "").upper() or "MESSAGE"
            created = str(item.get("created_at") or "")
            lines.append(f"[{created}] {role}" if created else role)
            lines.append(str(item.get("content") or ""))
            names = _attachment_names(item)
            if names:
                lines.append(f"Attachments: {names}")
            lines.append("")
        return "\n".join(lines).encode("utf-8")

    def _docx(self, title: str, messages: list[dict[str, Any]]) -> bytes:
        from docx import Document

        doc = Document()
        doc.add_heading(str(title or "Cuộc trò chuyện"), level=1)
        for item in messages:
            role = "Bạn" if item.get("role") == "user" else "Luyện"
            created = str(item.get("created_at") or "")
            heading = f"{role} · {created}" if created else role
            doc.add_heading(heading, level=2)
            for paragraph in str(item.get("content") or "").splitlines() or [""]:
                doc.add_paragraph(paragraph)
            names = _attachment_names(item)
            if names:
                doc.add_paragraph(f"Tệp: {names}")
        buffer = BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    def _pdf(self, title: str, messages: list[dict[str, Any]]) -> bytes:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.enums import TA_LEFT
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        buffer = BytesIO()
        font_name = "Helvetica"
        candidates = [
            Path("C:/Windows/Fonts/arial.ttf"),
            Path("C:/Windows/Fonts/calibri.ttf"),
            Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
            Path("/usr/share/fonts/dejavu/DejaVuSans.ttf"),
        ]
        for candidate in candidates:
            if candidate.is_file():
                try:
                    pdfmetrics.registerFont(TTFont("ExportUnicode", str(candidate)))
                    font_name = "ExportUnicode"
                    break
                except Exception:
                    continue

        def safe_text(value: str) -> str:
            text = str(value or "")
            if font_name == "Helvetica":
                text = text.encode("latin-1", "replace").decode("latin-1")
            return html.escape(text).replace("\n", "<br/>")

        styles = getSampleStyleSheet()
        title_style = styles["Title"]
        body_style = styles["BodyText"]
        heading_style = styles["Heading2"]
        for style in (title_style, body_style, heading_style):
            style.fontName = font_name
            style.leading = max(style.leading, 14)
            style.alignment = TA_LEFT

        story = [Paragraph(safe_text(title or "Cuộc trò chuyện"), title_style), Spacer(1, 12)]
        for item in messages:
            role = "Bạn" if item.get("role") == "user" else "Luyện"
            created = str(item.get("created_at") or "")
            story.append(Paragraph(safe_text(f"{role} · {created}" if created else role), heading_style))
            story.append(Paragraph(safe_text(str(item.get("content") or "")), body_style))
            names = _attachment_names(item)
            if names:
                story.append(Paragraph(safe_text(f"Tệp: {names}"), body_style))
            story.append(Spacer(1, 10))

        doc = SimpleDocTemplate(buffer, pagesize=A4, title=str(title or "Export"), rightMargin=36, leftMargin=36, topMargin=42, bottomMargin=42)
        doc.build(story)
        return buffer.getvalue()

    def _xlsx(self, title: str, messages: list[dict[str, Any]]) -> bytes:
        from openpyxl import Workbook
        from openpyxl.styles import Font
        from openpyxl.utils import get_column_letter

        workbook = Workbook()
        sheet = workbook.active
        table = _first_markdown_table(messages)
        if table:
            sheet.title = "Data"
            for row in table:
                sheet.append(row)
            for cell in sheet[1]:
                cell.font = Font(bold=True)
            sheet.freeze_panes = "A2"
            for column_index in range(1, len(table[0]) + 1):
                longest = max(
                    len(str(sheet.cell(row=row_index, column=column_index).value or ""))
                    for row_index in range(1, sheet.max_row + 1)
                )
                sheet.column_dimensions[get_column_letter(column_index)].width = min(60, max(12, longest + 2))
        else:
            sheet.title = "Conversation"
            headers = ["Time", "Role", "Content", "Attachments"]
            sheet.append(headers)
            for cell in sheet[1]:
                cell.font = Font(bold=True)
            for item in messages:
                sheet.append([
                    str(item.get("created_at") or ""),
                    str(item.get("role") or ""),
                    str(item.get("content") or ""),
                    _attachment_names(item),
                ])
            sheet.freeze_panes = "A2"
            sheet.column_dimensions["A"].width = 24
            sheet.column_dimensions["B"].width = 12
            sheet.column_dimensions["C"].width = 80
            sheet.column_dimensions["D"].width = 40
        buffer = BytesIO()
        workbook.save(buffer)
        return buffer.getvalue()

    def _pptx(self, title: str, messages: list[dict[str, Any]]) -> bytes:
        from pptx import Presentation
        from pptx.util import Pt

        content = "\n\n".join(
            str(item.get("content") or "").strip()
            for item in messages
            if str(item.get("content") or "").strip()
        )
        if not content:
            raise ArtifactError("Không có nội dung PowerPoint để xuất.")

        slide_header = re.compile(r"(?im)^\s*(?:slide|trang)\s*\d+\s*[:\-]\s*(.+?)\s*$")
        matches = list(slide_header.finditer(content))
        sections: list[tuple[str, str]] = []
        if matches:
            for idx, match in enumerate(matches):
                start = match.end()
                end = matches[idx + 1].start() if idx + 1 < len(matches) else len(content)
                sections.append((match.group(1).strip(), content[start:end].strip()))
        else:
            chunks = [chunk.strip() for chunk in re.split(r"\n\s*\n", content) if chunk.strip()]
            for index, chunk in enumerate(chunks, start=1):
                lines = [line.strip() for line in chunk.splitlines() if line.strip()]
                if not lines:
                    continue
                first = re.sub(r"^#+\s*", "", lines[0]).strip()
                body_lines = lines[1:] if len(lines) > 1 else []
                sections.append((first or f"Slide {index}", "\n".join(body_lines)))

        if not sections:
            sections = [(str(title or "Presentation"), content)]

        prs = Presentation()
        for slide_title, body in sections:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_title or str(title or "Slide")
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            body_lines = [line.strip() for line in body.splitlines() if line.strip()] or [""]
            for index, line in enumerate(body_lines):
                cleaned = re.sub(r"^(?:[-*•]|\d+[.)])\s*", "", line).strip()
                paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                paragraph.text = cleaned
                paragraph.font.size = Pt(22)

        buffer = BytesIO()
        prs.save(buffer)
        return buffer.getvalue()
